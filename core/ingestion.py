"""
Slide Library Ingestion Service

Handles ingestion of multi-slide presentations into the slide library.
Extracts individual slides, generates descriptions, and stores them.
"""

import logging
from pathlib import Path
from typing import List
from datetime import datetime
import tempfile
import hashlib
import zipfile

from utils.load_and_merge import PPTXLoader, PPTXSlideManager
from utils.utils import normalize_presentation, extract_slide_notes
from models.vertex import vertexai_model
from models.voyage import voyage_embed
from utils.schemas import (
    SlideLibraryMetadata, 
    SlideMetadata, 
    StorageReference,
    SlideContent
)
from prompts import SLIDE_DESCRIPTION_SYSTEM_PROMPT, SLIDE_DESCRIPTION_USER_PROMPT

from core.storage import SlideStorageAdapter

logger = logging.getLogger(__name__)


class SlideIngestionService:
    """
    Service for ingesting presentations into the slide library.
    
    Workflow:
    1. Load multi-slide presentation
    2. Extract each slide to single-slide PPTX
    3. Generate description (user notes > LLM)
    4. Create metadata
    5. Generate embedding
    6. Store atomically (S3 + MongoDB + Qdrant)
    """
    
    def __init__(self, storage: SlideStorageAdapter):
        """
        Initialize ingestion service.
        
        Args:
            storage: Storage adapter for slide library
        """
        self.storage = storage
        print("SlideIngestionService initialized")
    
    async def ingest_presentation(
        self,
        pptx_path: str
    ) -> List[SlideLibraryMetadata]:
        """
        Ingest a multi-slide presentation into the slide library.
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            List of SlideLibraryMetadata for each ingested slide
        """
        print(f"Starting ingestion: {pptx_path}")
        
        # Load presentation
        loader = PPTXLoader(pptx_path)
        presentation = loader.get_presentation()
        slide_count = loader.get_slide_count()
        dimensions = loader.get_dimensions()
        
        print(f"Loaded presentation: {slide_count} slides, {dimensions}")
        
        # Normalize to extract content structure
        _, content_mapping = normalize_presentation(pptx_path)
        
        # Process each slide
        ingested_slides = []
        temp_dir = Path(tempfile.mkdtemp(prefix="slide_library_"))
        
        try:
            for slide_idx in range(slide_count):
                print(f"Processing slide {slide_idx + 1}/{slide_count}")
                
                try:
                    # Extract single slide
                    single_slide_path = await self._extract_single_slide(
                        loader,
                        slide_idx,
                        temp_dir
                    )
                    
                    # Calculate file hash for deduplication
                    file_hash = self._calculate_file_hash(single_slide_path)
                    print(f"Slide hash: {file_hash[:16]}...")
                    
                    # Check if slide already exists
                    existing_slide = await self.storage.slide_exists_by_hash(file_hash)
                    if existing_slide:
                        print(f"⏭️  Slide already exists (hash: {file_hash[:16]}...), skipping")
                        print(f"   Existing slide ID: {existing_slide.slide_id}")
                        print(f"   Description: {existing_slide.description[:100]}...")
                        ingested_slides.append(existing_slide)
                        continue
                    
                    # Generate description (user notes > LLM)
                    description = await self._generate_description(
                        loader,
                        slide_idx,
                        content_mapping
                    )
                    
                    # Create metadata
                    metadata = SlideLibraryMetadata(
                        file_hash=file_hash,
                        description=description,
                        dimensions={
                            "width": int(dimensions["width"]),
                            "height": int(dimensions["height"])
                        },
                        element_count=len(content_mapping.slides[slide_idx].content) if slide_idx < len(content_mapping.slides) else 0,
                        storage_ref=StorageReference(
                            s3_key="",  # Will be filled by storage
                            mongodb_id="",
                            qdrant_id=""
                        ),
                        source_presentation=Path(pptx_path).name,
                        slide_index=slide_idx
                    )
                    
                    # Generate embedding
                    embedding = await self._generate_embedding(description)
                    
                    # Store atomically
                    storage_ref = await self.storage.store_slide(
                        slide_pptx_path=single_slide_path,
                        metadata=metadata,
                        embedding=embedding
                    )
                    
                    # Update metadata with storage references
                    metadata.storage_ref = storage_ref
                    
                    ingested_slides.append(metadata)
                    print(f"✅ Ingested slide {slide_idx + 1}: {metadata.slide_id}")
                    
                except Exception as e:
                    print(f"Failed to ingest slide {slide_idx + 1}: {e}")
                    # Continue with next slide
                    continue
            
            print(f"Ingestion complete: {len(ingested_slides)}/{slide_count} slides")
            return ingested_slides
            
        finally:
            # Cleanup temp directory
            try:
                import shutil
                shutil.rmtree(temp_dir)
                print(f"Cleaned up temp directory: {temp_dir}")
            except Exception as e:
                print(f"Failed to cleanup temp directory: {e}")
    
    async def _extract_single_slide(
        self,
        loader: PPTXLoader,
        slide_idx: int,
        temp_dir: Path
    ) -> Path:
        """
        Extract a single slide to a standalone PPTX file.
        
        Strategy: Load the entire source presentation and delete all slides
        except the target one. This preserves ALL original settings perfectly.
        
        Args:
            loader: PPTXLoader instance
            slide_idx: Index of slide to extract (0-based)
            temp_dir: Temporary directory for output
            
        Returns:
            Path to single-slide PPTX file
        """
        from spire.presentation import Presentation
        
        # Load a fresh copy of the source presentation
        # This preserves ALL settings: dimensions, themes, masters, etc.
        new_prs = Presentation()
        new_prs.LoadFromFile(loader.pptx_path)
        
        total_slides = new_prs.Slides.Count
        
        # Delete all slides EXCEPT the target slide
        # Work backwards to avoid index shifting issues
        for i in range(total_slides - 1, -1, -1):
            if i != slide_idx:
                new_prs.Slides.RemoveAt(i)
        
        # Verify we have exactly 1 slide
        if new_prs.Slides.Count != 1:
            raise RuntimeError(f"Expected 1 slide after extraction, got {new_prs.Slides.Count}")
        
        # Save to temp file
        output_path = temp_dir / f"slide_{slide_idx + 1}.pptx"
        PPTXSlideManager.save_presentation(new_prs, str(output_path))
        
        # Dispose
        new_prs.Dispose()
        
        return output_path
    
    async def _generate_description(
        self,
        loader: PPTXLoader,
        slide_idx: int,
        content_mapping
    ) -> str:
        """
        Generate rich description for a slide.
        
        PRIORITY: User notes > LLM generation
        
        Args:
            loader: PPTXLoader instance
            slide_idx: Index of slide (0-based)
            content_mapping: Normalized content mapping
            
        Returns:
            Rich description string
        """
        # Load slide with python-pptx to check notes
        from pptx import Presentation as PPTXPresentation
        pptx_prs = PPTXPresentation(loader.pptx_path)
        slide = pptx_prs.slides[slide_idx]
        
        # Check for user notes (GROUND TRUTH)
        notes_text = extract_slide_notes(slide)
        if notes_text:
            print(f"Using user notes as description (ground truth)")
            return notes_text
        
        # No notes - generate with LLM using full slide structure
        print(f"No user notes found, generating description with LLM")
        
        # Get slide content structure
        slide_content = None
        if slide_idx < len(content_mapping.slides):
            slide_content = content_mapping.slides[slide_idx]
        
        if not slide_content:
            print(f"No content mapping found for slide {slide_idx}")
            return f"Slide {slide_idx + 1} from presentation"
        
        # Prepare slide structure for LLM (exclude font and actual content)
        slide_structure = {
            "slide": slide_content.slide,
            "metadata": {
                "width": slide_content.metadata.width,
                "height": slide_content.metadata.height
            },
            "content": {}
        }
        
        # Extract relevant metadata for each component
        for uuid, content_item in slide_content.content.items():
            slide_structure["content"][uuid] = {
                "content_type": content_item.content_type,
                "position": {
                    "x": content_item.position.x,
                    "y": content_item.position.y
                },
                "size": {
                    "width": content_item.size.width,
                    "height": content_item.size.height
                },
                "content_description": content_item.content_description
            }
        
        user_prompt = SLIDE_DESCRIPTION_USER_PROMPT(slide_structure)

        try:
            description = await vertexai_model(
                system=SLIDE_DESCRIPTION_SYSTEM_PROMPT,
                user=user_prompt,
                temperature=0.3
            )
            
            print(f"Generated description: {description[:100]}...")
            return description.strip()
            
        except Exception as e:
            print(f"LLM description generation failed: {e}")
            # Fallback to basic description
            return f"Slide {slide_idx + 1} from presentation"
    
    async def _generate_embedding(self, description: str) -> list[float]:
        """
        Generate Voyage embedding for description.
        
        Args:
            description: Description text
            
        Returns:
            1024-dim embedding vector
        """
        try:
            # Generate embedding using voyage_embed
            embeddings = await voyage_embed(
                content=[description],
                input_type="document",
                model="voyage-3-large"
            )
            
            embedding = embeddings[0]  # Get first (and only) embedding
            print(f"Generated embedding: {len(embedding)} dimensions")
            return embedding
            
        except Exception as e:
            print(f"Embedding generation failed: {e}")
            raise
    def _calculate_file_hash(self, file_path: Path) -> str:
        """
        Calculate SHA256 hash of the entire file content.

        Args:
            file_path: Path to the file

        Returns:
            Hexadecimal hash string
        """
        sha256_hash = hashlib.sha256()

        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)

        return sha256_hash.hexdigest()
