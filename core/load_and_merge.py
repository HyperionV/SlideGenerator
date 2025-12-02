"""
PPTX loader and slide-copy utilities built on Spire.Presentation.

Main pieces:
- PPTXLoader(path): loads a .pptx from disk and exposes Spire's Presentation via get_presentation(), plus helpers like get_slides(), get_slide(i), get_slide_count(), get_dimensions().
- PPTXSlideManager: static helpers to copy slides, sync dimensions, and save with post-processing:
  - copy_slide(source_prs, source_idx, target_prs, target_position=None)
  - copy_slide_with_template(source_prs, source_idx, template_prs, template_idx)
  - copy_presentation_dimensions(source_prs, target_prs)
  - save_presentation(prs, output_path)  # saves and strips eval / empty Google bullet shapes

Typical usage:
    loader = PPTXLoader("input.pptx")
    prs = loader.get_presentation()
    PPTXSlideManager.copy_slide(prs, 0, prs)  # copy slide 0 to same deck
    PPTXSlideManager.copy_slide_with_template(prs, 0, prs, 0)  # copy slide 0 to same deck with template
    PPTXSlideManager.save_presentation(prs, "output.pptx")
"""

import os
from typing import List, Optional, Dict, Any
from spire.presentation import Presentation, FileFormat, SlideOrienation
from spire.presentation.common import SizeF
from pptx import Presentation as PPTXPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging

logger = logging.getLogger(__name__)


class PPTXLoader:
    """
    Loader class for PowerPoint presentations using Spire.Presentation.
    Provides complete slide copying with automatic preservation of all elements.
    """
    
    def __init__(self, pptx_path: str):
        """
        Initialize the loader with a PPTX file path.
        
        Args:
            pptx_path: Path to the PowerPoint file
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        self.pptx_path = pptx_path
        self.presentation = None
        self._load()
    
    def _load(self):
        """Load the presentation from file."""
        try:
            self.presentation = Presentation()
            self.presentation.LoadFromFile(self.pptx_path)
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise
    
    def get_presentation(self) -> Presentation:
        """
        Get the loaded presentation object.
        
        Returns:
            Presentation object from Spire.Presentation
        """
        return self.presentation
    
    def get_slides(self) -> List:
        """
        Get all slides from the presentation.
        
        Returns:
            List of Slide objects
        """
        return [self.presentation.Slides[i] for i in range(self.presentation.Slides.Count)]
    
    def get_slide(self, index: int) -> Optional[Any]:
        """
        Get a specific slide by index.
        
        Args:
            index: 0-based index of the slide
            
        Returns:
            Slide object or None if index is out of range
        """
        if 0 <= index < self.presentation.Slides.Count:
            return self.presentation.Slides[index]
        logger.warning(f"Slide index {index} out of range (0-{self.presentation.Slides.Count-1})")
        return None
    
    def get_slide_count(self) -> int:
        """
        Get the total number of slides.
        
        Returns:
            Number of slides in the presentation
        """
        return self.presentation.Slides.Count
    
    def get_slide_info(self, slide_index: int) -> Dict[str, Any]:
        """
        Get information about a specific slide.
        
        Args:
            slide_index: 0-based index of the slide
            
        Returns:
            Dictionary with slide information
        """
        slide = self.get_slide(slide_index)
        if slide is None:
            return {}
        
        info = {
            "index": slide_index,
            "shapes_count": slide.Shapes.Count,
            "has_notes": slide.NotesSlide is not None,
            "layout_name": "N/A"
        }
        
        return info
    
    def list_all_slides_info(self) -> List[Dict[str, Any]]:
        """
        Get information about all slides in the presentation.
        
        Returns:
            List of dictionaries with slide information
        """
        return [self.get_slide_info(i) for i in range(self.get_slide_count())]
    
    def get_dimensions(self) -> Dict[str, Any]:
        """
        Get presentation dimensions and orientation.
        
        Returns:
            Dictionary with width, height, and orientation
        """
        size = self.presentation.SlideSize.Size
        return {
            "width": size.Width,
            "height": size.Height,
            "orientation": "Portrait" if self.presentation.SlideSize.Orientation == SlideOrienation.Portrait else "Landscape"
        }
    
    def dispose(self):
        """Dispose of the presentation object to free resources."""
        if self.presentation is not None:
            self.presentation.Dispose()
            self.presentation = None


class PPTXSlideManager:
    """
    Manager class for slide operations across presentations using Spire.Presentation.
    Handles copying slides between presentations with complete element preservation.
    """
    
    @staticmethod
    def remove_all_shapes(slide: Any):
        """Remove all shapes from a slide."""
        for i in range(slide.Shapes.Count - 1, -1, -1):
            slide.Shapes.RemoveAt(i)
    
    @staticmethod
    def copy_slide_with_template(source_prs: Presentation, source_slide_index: int, 
                                  template_prs: Presentation, template_slide_index: int) -> Any:
        """Copy a slide from source to template presentation."""
        if source_slide_index < 0 or source_slide_index >= source_prs.Slides.Count:
            raise IndexError(f"Source slide index {source_slide_index} out of range (0-{source_prs.Slides.Count-1})")
        
        if template_slide_index < 0 or template_slide_index >= template_prs.Slides.Count:
            raise IndexError(f"Template slide index {template_slide_index} out of range (0-{template_prs.Slides.Count-1})")
        
        source_slide = source_prs.Slides[source_slide_index]
        template_slide = template_prs.Slides[template_slide_index]
        
        template_prs.Slides.AppendBySlide(source_slide)
        imported_slide = template_prs.Slides[template_prs.Slides.Count - 1]

        try:
            imported_slide.Layout = template_slide.Layout
        except Exception as e:
            logger.warning(f"Could not apply template layout: {e}")
        
        return imported_slide
    
    @staticmethod
    def copy_slide(source_prs: Presentation, source_slide_index: int, 
                   target_prs: Presentation, target_position: Optional[int] = None) -> Any:
        """Copy a slide from source presentation to target presentation."""
        if source_slide_index < 0 or source_slide_index >= source_prs.Slides.Count:
            raise IndexError(f"Source slide index {source_slide_index} out of range (0-{source_prs.Slides.Count-1})")
        
        source_slide = source_prs.Slides[source_slide_index]

        if target_position is not None:
            target_prs.Slides.Insert(target_position, source_slide)
            target_slide = target_prs.Slides[target_position]
        else:
            target_prs.Slides.AppendBySlide(source_slide)
            target_slide = target_prs.Slides[target_prs.Slides.Count - 1]

        return target_slide
    
    @staticmethod
    def copy_presentation_dimensions(source_prs: Presentation, target_prs: Presentation):
        try:
            target_prs.SlideSize.Size = SizeF(
                source_prs.SlideSize.Size.Width,
                source_prs.SlideSize.Size.Height
            )

            target_prs.SlideSize.Orientation = source_prs.SlideSize.Orientation

            target_prs.SlideSize.Type = source_prs.SlideSize.Type
        except Exception as e:
            logger.warning(f"Could not copy dimensions: {e}")
    
    @staticmethod
    def save_presentation(presentation: Presentation, output_path: str):
        try:
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)

            presentation.SaveToFile(output_path, FileFormat.Pptx2016)

            PPTXSlideManager.post_processing(output_path)
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise
    
    @staticmethod
    def post_processing(pptx_path: str):
        try:
            prs = PPTXPresentation(pptx_path)

            for slide_idx, slide in enumerate(prs.slides):
                shapes_to_remove = []

                for shape in slide.shapes:
                    if PPTXSlideManager._contains_evaluation_warning(shape) or PPTXSlideManager._is_empty_google_shape(shape):
                        shapes_to_remove.append(shape)
                
                for shape in shapes_to_remove:
                    sp = shape._sp
                    sp.getparent().remove(sp)

            prs.save(pptx_path)
        except Exception as e:
            logger.error(f"Failed during post-processing: {e}")
            raise
    
    @staticmethod
    def _contains_evaluation_warning(shape) -> bool:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            if "Evaluation Warning" in shape.text:
                return True

        if hasattr(shape, 'text'):
            if "Evaluation Warning" in shape.text:
                return True

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if PPTXSlideManager._contains_evaluation_warning(sub_shape):
                    return True
        
        return False
    
    @staticmethod
    def _is_empty_google_shape(shape) -> bool:
        if not hasattr(shape, 'name'):
            return False
        
        if shape.name != "Google Shape;18;p4":
            return False
        
        if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
            return False

        paragraphs = shape.text_frame.paragraphs
        if len(paragraphs) != 1:
            return False
        
        paragraph = paragraphs[0]
        
        paragraph_text = paragraph.text.strip() if paragraph.text else ""
        if paragraph_text:
            return False

        has_bullet = False

        if hasattr(paragraph, 'level'):
            try:
                if paragraph.level is not None and paragraph.level >= 0:
                    has_bullet = True
            except:
                pass
        
        if not has_bullet and hasattr(paragraph, '_element'):
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    if (hasattr(pPr, 'buFont') or 
                        hasattr(pPr, 'buChar') or 
                        hasattr(pPr, 'buAutoNum') or
                        hasattr(pPr, 'buBlip')):
                        has_bullet = True
            except:
                pass
        
        return has_bullet


def load_pptx(pptx_path: str) -> PPTXLoader:
    """
    Convenience function to load a PPTX file.
    
    Args:
        pptx_path: Path to the PowerPoint file
        
    Returns:
        PPTXLoader instance
    """
    return PPTXLoader(pptx_path)
