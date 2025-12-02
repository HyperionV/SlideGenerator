import os
import uuid
import json
from pptx import Presentation
from typing import Dict, List, Union, Any
from dataclasses import dataclass, asdict


@dataclass
class Position:
    x: int
    y: int


@dataclass
class Size:
    width: int
    height: int


@dataclass
class Font:
    size: int
    bold: bool
    italic: bool = False
    underline: bool = False
    name: str = "Arial"


@dataclass
class ContentItem:
    original_content: Union[str, List[List[str]]]
    position: Position
    size: Size
    alignment: str  # "left", "center", "right", "justify"
    font: Font


@dataclass
class SlideMetadata:
    width: int
    height: int


@dataclass
class SlideContent:
    slide: int
    metadata: SlideMetadata
    content: Dict[str, ContentItem]  # UUID -> ContentItem


@dataclass
class PresentationMapping:
    slides: List[SlideContent]

def replace_text_with_uuid(pptx_path):
    # Load the presentation
    presentation = Presentation(pptx_path)

    # Structure to map original content with UUIDs
    content_mapping = PresentationMapping(slides=[])

    # Get presentation dimensions (convert from EMUs to pixels)
    slide_width_px = presentation.slide_width // 12700  # EMU to pixel conversion
    slide_height_px = presentation.slide_height // 12700

    # Iterate through all slides
    for slide_idx, slide in enumerate(presentation.slides, 1):
        slide_metadata = SlideMetadata(width=slide_width_px, height=slide_height_px)
        slide_content = SlideContent(slide=slide_idx, metadata=slide_metadata, content={})

        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            is_textbox = shape.has_text_frame
            is_table = shape.has_table

            if is_textbox or is_table:
                should_replace = False
                original_content = ""
                content_type = ""

                # Get position and size (convert from EMUs to pixels)
                position = Position(
                    x=shape.left // 12700,
                    y=shape.top // 12700
                )
                size = Size(
                    width=shape.width // 12700,
                    height=shape.height // 12700
                )

                if is_textbox:
                    textbox_text = shape.text_frame.text.strip()
                    should_replace = textbox_text and not textbox_text.replace('.', '').replace(',', '').isdigit()
                    original_content = textbox_text
                    content_type = "TEXT"

                elif is_table:
                    table = shape.table
                    table_content = []
                    for row in table.rows:
                        row_content = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_content.append(cell_text)
                            if cell_text and not cell_text.replace('.', '').replace(',', '').isdigit():
                                should_replace = True
                        table_content.append(row_content)
                    original_content = table_content
                    content_type = "TABLE"

                if should_replace:
                    text_uuid = str(uuid.uuid4())

                    # Extract alignment and font information
                    alignment = "left"  # default
                    font_info = Font(size=12, bold=False)  # defaults

                    if is_textbox and shape.text_frame.paragraphs:
                        paragraph = shape.text_frame.paragraphs[0]
                        # Get alignment
                        alignment_map = {
                            0: "left",
                            1: "center",
                            2: "right",
                            3: "justify"
                        }
                        alignment = alignment_map.get(paragraph.alignment, "left")

                        # Get font info from first run
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            font_size = run.font.size
                            if font_size:
                                font_size = font_size // 12700  # Convert from EMUs to points
                            else:
                                font_size = 12  # default

                            font_info = Font(
                                size=font_size,
                                bold=bool(run.font.bold),
                                italic=bool(run.font.italic),
                                underline=bool(run.font.underline),
                                name=run.font.name or "Arial"
                            )

                    # Add to slide content mapping
                    slide_content.content[text_uuid] = ContentItem(
                        original_content=original_content,
                        position=position,
                        size=size,
                        alignment=alignment,
                        font=font_info
                    )

                    if is_textbox:
                        text_frame = shape.text_frame
                        while len(text_frame.paragraphs) > 1:
                            paragraph_to_remove = text_frame.paragraphs[-1]
                            p = paragraph_to_remove._element
                            p.getparent().remove(p)
                            paragraph_to_remove._p = paragraph_to_remove._element = None

                        if text_frame.paragraphs:
                            paragraph = text_frame.paragraphs[0]
                            original_run = paragraph.runs[0] if paragraph.runs else None

                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = text_uuid

                            if original_run:
                                run.font.size = original_run.font.size
                                run.font.bold = original_run.font.bold
                                run.font.italic = original_run.font.italic
                                run.font.underline = original_run.font.underline
                                try:
                                    run.font.color.rgb = original_run.font.color.rgb
                                except AttributeError:
                                    from pptx.dml.color import RGBColor
                                    run.font.color.rgb = RGBColor(0, 0, 0)
                                run.font.name = original_run.font.name

                    elif is_table:
                        table = shape.table
                        num_rows = len(table.rows)
                        num_cols = len(table.columns)

                        # Remove all rows except the first
                        for _ in range(num_rows - 1):
                            table._tbl.remove(table._tbl.tr_lst[-1])

                        # Remove all columns except the first
                        for _ in range(num_cols - 1):
                            for row in table.rows:
                                row._tr.remove(row._tr.tc_lst[-1])

                        # Set UUID in the single remaining cell
                        table.rows[0].cells[0].text = text_uuid

        # Add slide content to mapping if it has any content
        if slide_content.content:
            content_mapping.slides.append(slide_content)

    return presentation, content_mapping

def export_slide_structure(content_mapping):
    """
    Export the content mapping to the specified JSON structure.

    Args:
        content_mapping: PresentationMapping object

    Returns:
        List of slide dictionaries in the requested format
    """
    slides_list = []

    for slide_content in content_mapping.slides:
        slide_dict = {
            "slide": slide_content.slide,
            "metadata": {
                "width": slide_content.metadata.width,
                "height": slide_content.metadata.height
            },
            "content": {}
        }

        for uuid_key, content_item in slide_content.content.items():
            slide_dict["content"][uuid_key] = {
                "original_content": content_item.original_content,
                "position": {
                    "x": content_item.position.x,
                    "y": content_item.position.y
                },
                "size": {
                    "width": content_item.size.width,
                    "height": content_item.size.height
                },
                "alignment": content_item.alignment,
                "font": {
                    "size": content_item.font.size,
                    "bold": content_item.font.bold,
                    "italic": content_item.font.italic,
                    "underline": content_item.font.underline,
                    "name": content_item.font.name
                }
            }

        slides_list.append(slide_dict)

    return slides_list

def update_content(presentation_path, structure_path, output_path=None):
    """
    Update slide content based on processed structure.

    Args:
        presentation_path: Path to the PowerPoint file to update
        structure_path: Path to the JSON file containing the update structure
        output_path: Path to save the updated presentation (optional, defaults to modified input path)
    """
    import json

    # Load the presentation
    presentation = Presentation(presentation_path)

    # Load the structure from JSON file
    with open(structure_path, 'r', encoding='utf-8') as f:
        structure_data = json.load(f)

    # Extract content mapping from structure
    content_mapping = structure_data.get('content', {})

    # Process all slides
    for slide in presentation.slides:
        # Process each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                current_text = shape.text_frame.text.strip()
                if current_text in content_mapping:
                    # Replace with the content from structure
                    new_content = content_mapping[current_text]['content']

                    # Clear existing paragraphs and add new content
                    text_frame = shape.text_frame
                    while len(text_frame.paragraphs) > 1:
                        paragraph_to_remove = text_frame.paragraphs[-1]
                        p = paragraph_to_remove._element
                        p.getparent().remove(p)
                        paragraph_to_remove._p = paragraph_to_remove._element = None

                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        original_run = paragraph.runs[0] if paragraph.runs else None

                        paragraph.clear()
                        run = paragraph.add_run()
                        run.text = new_content

                        # Preserve original formatting if available
                        if original_run:
                            run.font.size = original_run.font.size
                            run.font.bold = original_run.font.bold
                            run.font.italic = original_run.font.italic
                            run.font.underline = original_run.font.underline
                            try:
                                run.font.color.rgb = original_run.font.color.rgb
                            except AttributeError:
                                from pptx.dml.color import RGBColor
                                run.font.color.rgb = RGBColor(0, 0, 0)
                            run.font.name = original_run.font.name

    # Determine output path
    if output_path is None:
        base_name = os.path.splitext(presentation_path)[0]
        output_path = f"{base_name}_updated.pptx"

    # Save the updated presentation
    presentation.save(output_path)
    print(f"Updated presentation saved to {output_path}")

    return output_path

def main():
    # Input file path
    input_path = "input/test.pptx"

    # Output file path for UUID presentation
    uuid_output_path = "output/test_uuid.pptx"

    # Process the presentation and get content mapping
    presentation, content_mapping = replace_text_with_uuid(input_path)

    # Save the modified presentation with UUIDs
    presentation.save(uuid_output_path)
    print(f"Processed presentation saved to {uuid_output_path}")

    # Export the slide structure to JSON
    slide_structure = export_slide_structure(content_mapping)

    # Save the structure to a JSON file
    structure_output_path = "output/slide_structure.json"
    with open(structure_output_path, 'w', encoding='utf-8') as f:
        json.dump(slide_structure, f, indent=2, ensure_ascii=False)

    print(f"Slide structure exported to {structure_output_path}")

    # Print the structure for verification
    print("\n=== SLIDE STRUCTURE ===")
    print(json.dumps(slide_structure, indent=2, ensure_ascii=False))

    # # Uncomment to test content update
    # structure_path = "processed_slide.json"
    # presentation_path = "input/test_input.pptx"
    # output_path = "output/test_updated.pptx"
    # update_content(presentation_path, structure_path, output_path=output_path)

if __name__ == "__main__":
    main()
