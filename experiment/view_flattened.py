from pptx import Presentation

def display_element(element, level=0):
    indent = "  " * level

    if hasattr(element, 'text') and not hasattr(element, 'cells'):  # Text element
        text = element.text.strip()
        if text:
            truncated = text[:50] + ('...' if len(text) > 50 else '')
            print(f"{indent}- text: {truncated}")

    elif hasattr(element, 'rows'):  # Table
        print(f"{indent}- table")
        for row_idx, row in enumerate(element.rows):
            print(f"{indent}  - row {row_idx + 1}")
            for cell_idx, cell in enumerate(row.cells):
                display_element(cell, level + 2)

    elif hasattr(element, 'chart_title'):  # Chart
        print(f"{indent}- chart")

    elif hasattr(element, 'paragraphs'):  # Text frame
        for paragraph in element.paragraphs:
            display_element(paragraph, level)

    elif hasattr(element, 'runs'):  # Paragraph
        for run in element.runs:
            display_element(run, level)

def group_elements_by_vertical_adjacency(shapes, vertical_threshold=50000, horizontal_threshold_percent=10, slide_width=None):  # threshold in EMUs (~0.5 inches)
    """Group elements by vertical adjacency and horizontal alignment, then sort groups left to right"""
    if not shapes:
        return []

    # Calculate horizontal threshold based on slide width
    if slide_width is None:
        slide_width = 9144000  # Default 10 inches in EMUs
    horizontal_threshold = slide_width * (horizontal_threshold_percent / 100.0)

    # Sort shapes by top position first
    sorted_shapes = sorted(shapes, key=lambda s: s.top)

    # Group shapes by vertical proximity AND horizontal alignment
    vertical_groups = []
    current_group = [sorted_shapes[0]]

    for shape in sorted_shapes[1:]:
        # Check if shape is vertically adjacent to the last shape in current group
        last_shape = current_group[-1]
        vertical_distance = shape.top - (last_shape.top + last_shape.height)

        # Check if shape is horizontally aligned (within threshold of group's horizontal range)
        group_left = min(s.left for s in current_group)
        group_right = max(s.left + s.width for s in current_group)
        shape_left = shape.left
        shape_right = shape.left + shape.width

        # Check if there's horizontal overlap or proximity
        horizontal_overlap = max(0, min(group_right, shape_right) - max(group_left, shape_left))
        horizontal_distance = min(abs(shape_left - group_right), abs(shape_right - group_left))

        # Consider horizontally aligned if there's overlap OR distance is within threshold
        is_horizontally_aligned = horizontal_overlap > 0 or horizontal_distance <= horizontal_threshold

        if vertical_distance <= vertical_threshold and is_horizontally_aligned:
            # Shape is both vertically adjacent and horizontally aligned, add to current group
            current_group.append(shape)
        else:
            # Shape is not sufficiently aligned, start new group
            vertical_groups.append(current_group)
            current_group = [shape]

    # Add the last group
    if current_group:
        vertical_groups.append(current_group)

    # Sort each vertical group horizontally (left to right)
    for group in vertical_groups:
        group.sort(key=lambda s: s.left)

    # Sort the vertical groups themselves by their leftmost element's position
    vertical_groups.sort(key=lambda g: g[0].left)

    return vertical_groups

def display_flattened_structure(pptx_path):
    # Load the presentation
    presentation = Presentation(pptx_path)

    for slide_idx, slide in enumerate(presentation.slides, 1):
        print(f"slide {slide_idx}:")

        # Get slide dimensions for bounds checking
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Filter shapes that are within bounds
        valid_shapes = []
        for shape in slide.shapes:
            # Skip shapes that are entirely out of bounds
            if (shape.left + shape.width < 0 or
                shape.top + shape.height < 0 or
                shape.left > slide_width or
                shape.top > slide_height):
                continue
            valid_shapes.append(shape)

        # Group elements by vertical adjacency and horizontal alignment, then sort groups left to right
        vertical_groups = group_elements_by_vertical_adjacency(valid_shapes, horizontal_threshold_percent=10, slide_width=slide_width)

        for group_idx, group in enumerate(vertical_groups):
            print(f"  - group {group_idx + 1}:")

            for shape in group:
                if hasattr(shape, "text_frame"):
                    display_element(shape.text_frame, level=2)

                elif shape.shape_type == 19:  # TABLE shape type
                    try:
                        table = shape.table
                        display_element(table, level=2)
                    except:
                        pass

                elif shape.shape_type == 3:  # CHART shape type
                    display_element(shape, level=2)

        print()

def main():
    input_path = "input/test.pptx"

    display_flattened_structure(input_path)

if __name__ == "__main__":
    main()
