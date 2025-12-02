import os
import uuid
import json
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Any, Tuple
from schemas import (
    Position, Size, Font, ContentItem, SlideMetadata,
    SlideContent, PresentationMapping, ChartMetadata, ChartSeries,
    CellStyle, TableStyleInfo
)


def get_shape_alt_text(shape) -> str:
    """
    Get alt text from a shape.
    
    Args:
        shape: PPTX shape object
    
    Returns:
        str: Alt text if available, empty string otherwise
    """
    try:
        if hasattr(shape, 'alt_text') and shape.alt_text:
            return shape.alt_text.strip()
        elif hasattr(shape, '_element') and hasattr(shape._element, '_nvXxPr'):
            return shape._element._nvXxPr.cNvPr.attrib.get("descr", "").strip()
    except Exception:
        pass
    return ""


def set_shape_alt_text(shape, text: str) -> bool:
    """
    Set alt text for a shape.
    
    Args:
        shape: PPTX shape object
        text: Text to set as alt text
    
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        if hasattr(shape, 'alt_text'):
            shape.alt_text = text
            return True
        elif hasattr(shape, '_element') and hasattr(shape._element, '_nvXxPr'):
            shape._element._nvXxPr.cNvPr.attrib["descr"] = text
            return True
    except Exception as e:
        print(f"Warning: Failed to set alt text: {e}")
    return False


def break_external_chart_links(pptx_path: str) -> str:
    """
    Break external links in charts by removing externalData references and relationships.
    This converts linked charts to embedded charts that python-pptx can work with.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        Path to the modified PowerPoint file (may be the same or a temp file)
    """
    import tempfile
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
    os.close(temp_fd)

    try:
        # Extract PPTX (which is a ZIP file)
        with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
            with zipfile.ZipFile(temp_path, 'w') as temp_zip:
                for item in pptx_zip.filelist:
                    # Read the file content
                    content = pptx_zip.read(item.filename)

                    # Check if this is a chart XML file
                    if item.filename.startswith('ppt/charts/') and item.filename.endswith('.xml'):
                        # print(f"  Processing chart XML: {item.filename}")
                        # Parse and modify chart XML
                        modified_content = _modify_chart_xml(content)
                        temp_zip.writestr(item.filename, modified_content)
                    elif item.filename.endswith('.rels') and 'ppt/' in item.filename:
                        # print(f"  Processing relationships XML: {item.filename}")
                        modified_content = _modify_relationships_xml(content)
                        temp_zip.writestr(item.filename, modified_content)
                    else:
                        # Copy other files as-is
                        temp_zip.writestr(item.filename, content)

        return temp_path

    except Exception as e:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise e


def _modify_chart_xml(xml_content: bytes) -> bytes:
    """
    Remove externalData elements from chart XML to convert linked charts to embedded.
    """
    try:
        # Convert to string for easier manipulation
        xml_str = xml_content.decode('utf-8')

        # Look for externalData elements and remove them
        import re

        # Pattern to match <c:externalData r:id="rIdX">...</c:externalData>
        pattern = r'<c:externalData[^>]*r:id="rId\d+"[^>]*>.*?</c:externalData>'
        matches = re.findall(pattern, xml_str, re.DOTALL)

        external_data_count = 0
        for match in matches:
            xml_str = xml_str.replace(match, '', 1)  # Remove only first occurrence
            external_data_count += 1
            # print(f"    Removed externalData element: {match}")

        if external_data_count > 0:
            print(f"  Removed {external_data_count} externalData elements from chart")

        # Convert back to bytes
        return xml_str.encode('utf-8')

    except Exception as e:
        # If modification fails, return original content
        print(f"  Warning: Failed to modify chart XML: {e}")
        return xml_content


def _modify_relationships_xml(xml_content: bytes) -> bytes:
    """
    Remove external relationships from relationships XML.
    """
    try:
        # Parse XML
        root = ET.fromstring(xml_content)

        # Handle namespace - relationships XML uses a default namespace
        ns = {'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'}
        # Try both with and without namespace
        relationships = root.findall('.//Relationship') + root.findall('.//rel:Relationship', ns)

        # Find and remove external relationships
        relationships_to_remove = []
        external_count = 0
        for rel in relationships:
            target_mode = rel.get('TargetMode')
            rel_id = rel.get('Id')
            rel_type = rel.get('Type')
            target = rel.get('Target')
            # print(f"    Checking relationship: Id={rel_id}, Target={target}, TargetMode={target_mode}")
            if target_mode == 'External':
                relationships_to_remove.append(rel)
                external_count += 1

        # Remove the external relationships
        for rel in relationships_to_remove:
            root.remove(rel)

        if external_count > 0:
            print(f"  Removed {external_count} external relationships")

        # Convert back to bytes
        return ET.tostring(root, encoding='utf-8', xml_declaration=True)

    except Exception as e:
        # If modification fails, return original content
        print(f"  Warning: Failed to modify relationships XML: {e}")
        return xml_content


def extract_chart_metadata(chart) -> ChartMetadata:
    """
    Extract chart series names, values, and categories from a PowerPoint chart.

    Args:
        chart: The chart object from python-pptx

    Returns:
        ChartMetadata object containing series names, values, and categories
    """
    try:
        # Extract series data (names and values)
        series_list = []
        if hasattr(chart, 'series'):
            for series in chart.series:
                series_name = ""
                if hasattr(series, 'name') and series.name:
                    series_name = str(series.name)

                values = []
                if hasattr(series, 'values'):
                    values = list(series.values) if series.values else []

                series_obj = ChartSeries(
                    name=series_name,
                    values=values
                )
                series_list.append(series_obj)

        # Extract categories
        categories = []
        if hasattr(chart, 'categories') and chart.categories:
            categories = [str(cat) for cat in chart.categories]

        return ChartMetadata(series=series_list, categories=categories)

    except Exception as e:
        # Return empty metadata if extraction fails
        return ChartMetadata(series=[])


def extract_slide_notes(slide) -> str:
    """
    Extract notes text from a PowerPoint slide.

    Args:
        slide: PPTX slide object

    Returns:
        str: Notes text if available, empty string otherwise
    """
    try:
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                return notes_text
    except Exception as e:
        # Silently handle any errors in notes extraction
        pass
    return ""


def extract_cell_style(cell) -> CellStyle:
    """
    Extract comprehensive styling information from a PowerPoint table cell.
    
    Captures font properties (size, bold, italic, name, color) and cell fill color
    from the first text run in the cell. Used to preserve table formatting when
    modifying table structure or content.
    
    Args:
        cell: python-pptx table cell object
    
    Returns:
        CellStyle: Object containing extracted font and fill styling properties.
                   Returns default CellStyle if extraction fails or cell is empty.
    """
    style = CellStyle()
    
    try:
        if cell.text_frame and cell.text_frame.paragraphs:
            paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                run = paragraph.runs[0]
                font = run.font
                
                if font.size:
                    style.font_size = font.size // 12700
                
                style.font_bold = bool(font.bold)
                style.font_italic = bool(font.italic)
                style.font_name = font.name or "Arial"
                
                try:
                    if font.color and font.color.type and hasattr(font.color, 'rgb') and font.color.rgb:
                        rgb = font.color.rgb
                        style.font_color_rgb = (rgb[0], rgb[1], rgb[2])
                except:
                    pass
        
        try:
            fill = cell.fill
            if fill and fill.type:
                style.has_fill = True
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                        rgb = fill.fore_color.rgb
                        style.fill_color_rgb = (rgb[0], rgb[1], rgb[2])
        except:
            pass
    
    except Exception:
        pass
    
    return style


def extract_table_styling(table) -> TableStyleInfo:
    """
    Extract header and content cell styling from a PowerPoint table.
    
    Implements a two-cell sampling strategy:
    - Header style: Extracted from top-left cell [0,0]
    - Content style: Extracted from bottom-right cell [n-1,m-1]
    
    This approach assumes the first row contains headers and subsequent rows
    contain content, which is the standard table structure in presentations.
    
    Args:
        table: python-pptx table object
    
    Returns:
        TableStyleInfo: Object containing both header_style and content_style.
                        For single-row tables, uses rightmost cell for content style.
    """
    num_rows = len(table.rows)
    num_cols = len(table.columns)
    
    header_cell = table.rows[0].cells[0]
    header_style = extract_cell_style(header_cell)
    
    if num_rows > 1:
        content_cell = table.rows[num_rows - 1].cells[num_cols - 1]
    else:
        content_cell = table.rows[0].cells[num_cols - 1] if num_cols > 1 else table.rows[0].cells[0]
    
    content_style = extract_cell_style(content_cell)
    
    return TableStyleInfo(
        header_style=header_style,
        content_style=content_style
    )


def normalize_presentation(pptx_path: str) -> Tuple[Presentation, PresentationMapping]:
    """
    Load presentation and normalize the content, creating content mapping.
    Check slide notes and use them as purpose if available.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        Tuple of (modified_presentation, content_mapping)
    """
    # First, break any external chart links to make them work with python-pptx
    try:
        modified_pptx_path = break_external_chart_links(pptx_path)
        print(f"✅ Successfully processed chart links in: {pptx_path}")
    except Exception as e:
        print(f"⚠️  Failed to break chart links, proceeding with original file: {e}")
        modified_pptx_path = pptx_path

    # Load the presentation (using modified file if link-breaking succeeded)
    presentation = Presentation(modified_pptx_path)

    # Structure to map original content
    content_mapping = PresentationMapping(slides=[])

    slide_width_px = presentation.slide_width // 12700
    slide_height_px = presentation.slide_height // 12700

    for slide_idx, slide in enumerate(presentation.slides, 1):
        slide_metadata = SlideMetadata(width=slide_width_px, height=slide_height_px)

        # Check for slide notes and use as purpose if available
        slide_notes = extract_slide_notes(slide)
        slide_purpose = slide_notes if slide_notes else ""

        slide_content = SlideContent(slide=slide_idx, metadata=slide_metadata, content={}, purpose=slide_purpose)

        for shape in slide.shapes:
            is_textbox = shape.has_text_frame
            is_table = shape.has_table
            is_chart = shape.has_chart

            if is_textbox or is_table or is_chart:
                should_replace = False
                original_content = ""
                content_type = ""

                position = Position(
                    x=shape.left // 12700,
                    y=shape.top // 12700
                )
                size = Size(
                    width=shape.width // 12700,
                    height=shape.height // 12700
                )

                if is_textbox:
                    textbox_text = shape.text_frame.text.strip()
                    should_replace = textbox_text and not textbox_text.replace('.', '').replace(',', '').isdigit()
                    original_content = textbox_text
                    content_type = "TEXT"

                elif is_table:
                    table = shape.table
                    num_rows = len(table.rows)
                    num_cols = len(table.columns)
                    total_cells = num_rows * num_cols
                    
                    # Check if this is a 1-cell table
                    if total_cells == 1:
                        # Treat as TEXT component
                        cell_text = table.rows[0].cells[0].text.strip()
                        should_replace = bool(cell_text and not cell_text.replace('.', '').replace(',', '').isdigit())
                        original_content = cell_text
                        content_type = "TEXT"  # Mark as TEXT for generation
                        is_single_cell_table = True
                        table_style_info = None
                    else:
                        # Multi-cell table - keep as TABLE
                        table_content = []
                        for row in table.rows:
                            row_content = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_content.append(cell_text)
                                should_replace = True
                            table_content.append(row_content)
                        original_content = table_content
                        content_type = "TABLE"
                        is_single_cell_table = False
                        
                        # Extract table styling for multi-cell tables
                        table_style_info = extract_table_styling(table)

                elif is_chart:
                    chart_metadata = extract_chart_metadata(shape.chart)
                    original_content = chart_metadata
                    content_type = "CHART"
                    should_replace = True
                    is_single_cell_table = False
                else:
                    is_single_cell_table = False

                if should_replace:
                    text_uuid = str(uuid.uuid4())
                    
                    # Add "-cell" suffix for 1-cell tables
                    if is_table and is_single_cell_table:
                        text_uuid = f"{text_uuid}-cell"

                    font_info = Font(size=9, bold=False)

                    if is_textbox and shape.text_frame.paragraphs:
                        paragraph = shape.text_frame.paragraphs[0]

                        if paragraph.runs:
                            run = paragraph.runs[0]
                            font_size = run.font.size
                            if font_size:
                                font_size = font_size // 12700
                            else:
                                font_size = 9

                            font_info = Font(
                                size=font_size,
                                bold=bool(run.font.bold),
                                italic=bool(run.font.italic),
                                underline=bool(run.font.underline),
                                name=run.font.name or "Arial"
                            )
                    
                    elif is_table and is_single_cell_table:
                        # Get font from 1-cell table
                        cell = table.rows[0].cells[0]
                        if cell.text_frame and cell.text_frame.paragraphs:
                            paragraph = cell.text_frame.paragraphs[0]
                            if paragraph.runs:
                                run = paragraph.runs[0]
                                font_size = run.font.size
                                if font_size:
                                    font_size = font_size // 12700
                                else:
                                    font_size = 9
                                
                                font_info = Font(
                                    size=font_size,
                                    bold=bool(run.font.bold),
                                    italic=bool(run.font.italic),
                                    underline=bool(run.font.underline),
                                    name=run.font.name or "Arial"
                                )

                    # Create ContentItem with table_style for TABLE types
                    content_item_kwargs = {
                        'original_content': original_content,
                        'content_type': content_type,
                        'position': position,
                        'size': size,
                        'font': font_info
                    }
                    
                    # Add table_style only for multi-cell tables
                    if content_type == "TABLE" and not is_single_cell_table:
                        content_item_kwargs['table_style'] = table_style_info
                    
                    slide_content.content[text_uuid] = ContentItem(**content_item_kwargs)
                    
                    # Place UUID in alt_text instead of content
                    set_shape_alt_text(shape, text_uuid)


        if slide_content.content:
            content_mapping.slides.append(slide_content)

    return presentation, content_mapping


def export_slide_structure(content_mapping: PresentationMapping) -> List[Dict[str, Any]]:
    """
    Export the content mapping to the specified JSON structure.

    Args:
        content_mapping: PresentationMapping object

    Returns:
        List of slide dictionaries in the requested format
    """
    slides_list = []

    for slide_content in content_mapping.slides:
        slide_dict = {
            "slide": slide_content.slide,
            "purpose": slide_content.purpose,
            "content": {}
        }

        for uuid_key, content_item in slide_content.content.items():
            # Convert multi-line content for proper rendering
            display_content = content_item.content
            if content_item.content and isinstance(content_item.content, str):
                if '\n' in content_item.content:
                    # Split by newlines and create bullet list
                    lines = content_item.content.split('\n')
                    # Filter out empty lines
                    lines = [line.strip() for line in lines if line.strip()]
                    display_content = lines

            if content_item.content_type == "CHART" and hasattr(content_item.original_content, 'model_dump'):
                display_content = content_item.original_content.model_dump()
            elif content_item.content_type == "CHART":
                display_content = {
                    "series": getattr(content_item.original_content, 'series', [])
                }


            slide_dict["content"][uuid_key] = {
                "original_content": content_item.original_content.model_dump() if hasattr(content_item.original_content, 'model_dump') else content_item.original_content,
                "content_type": content_item.content_type,
                "position": {
                    "x": content_item.position.x,
                    "y": content_item.position.y
                },
                "size": {
                    "width": content_item.size.width,
                    "height": content_item.size.height
                },
                "font": {
                    "size": content_item.font.size,
                    "bold": content_item.font.bold,
                    "italic": content_item.font.italic,
                    "underline": content_item.font.underline,
                    "name": content_item.font.name
                },
                "content_description": content_item.content_description,
                "content": display_content,
            }
            
            # Add table_style for TABLE content types
            if content_item.table_style:
                slide_dict["content"][uuid_key]["table_style"] = content_item.table_style.model_dump()


        slides_list.append(slide_dict)

    return slides_list


def update_text_component(shape, content_item) -> bool:
    """
    Update a text component with new content while preserving formatting.

    Args:
        shape: PPTX shape object with text frame
        content_item: ContentItem containing the new content

    Returns:
        bool: True if update was successful
    """
    if not shape.has_text_frame or not content_item.content:
        return False

    text_frame = shape.text_frame
    new_content = content_item.content

    if not text_frame.paragraphs:
        return False

    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        return False

    uuid_run = paragraph.runs[0]

    if isinstance(new_content, list):
        while len(text_frame.paragraphs) > 1:
            paragraph_to_remove = text_frame.paragraphs[-1]
            p = paragraph_to_remove._element
            p.getparent().remove(p)
            paragraph_to_remove._p = paragraph_to_remove._element = None

        text_frame.paragraphs[0].clear()

        for i, line in enumerate(new_content):
            if i > 0:
                paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = f"• {line}" if line else ""

            # Copy formatting from original run
            _copy_font_formatting(run.font, uuid_run.font)
    else:
        uuid_run.text = str(new_content)

    return True


def update_single_cell_table(shape, content_item) -> bool:
    """
    Update a single-cell table with new content (treated as text).
    
    Args:
        shape: PPTX shape object with table
        content_item: ContentItem containing the new content
    
    Returns:
        bool: True if update was successful
    """
    if not shape.has_table or not content_item.content:
        return False
    
    table = shape.table
    if len(table.rows) != 1 or len(table.columns) != 1:
        return False
    
    cell = table.rows[0].cells[0]
    new_content = content_item.content
    
    if not cell.text_frame:
        return False
    
    text_frame = cell.text_frame
    
    # Get original formatting
    original_run = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        original_run = text_frame.paragraphs[0].runs[0]
    
    # Handle list content
    if isinstance(new_content, list):
        # Clear existing paragraphs
        while len(text_frame.paragraphs) > 1:
            paragraph_to_remove = text_frame.paragraphs[-1]
            p = paragraph_to_remove._element
            p.getparent().remove(p)
            paragraph_to_remove._p = paragraph_to_remove._element = None
        
        text_frame.paragraphs[0].clear()
        
        for i, line in enumerate(new_content):
            if i > 0:
                paragraph = text_frame.add_paragraph()
            else:
                paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = f"• {line}" if line else ""
            
            if original_run:
                _copy_font_formatting(run.font, original_run.font)
    else:
        # Single text content
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = str(new_content)
        
        if original_run:
            _copy_font_formatting(run.font, original_run.font)
    
    return True


def update_table_component(shape, content_item) -> bool:
    """
    Update a PowerPoint table with new data while preserving formatting and dimensions.
    
    Implements a three-phase update strategy:
    1. Structure Adjustment: Adds/removes rows and columns to match new data dimensions,
       then redistributes column widths to maintain original table width.
    2. Style Application: Applies header styling to row 0 and content styling to
       remaining rows using stored table_style information.
    3. Text Population: Sets cell text while preserving formatting through a
       store-clear-reapply pattern that handles both existing and newly added cells.
    
    The function handles tables with varying dimensions and ensures new cells inherit
    appropriate styling even when they have no initial formatting (empty cells from
    column/row addition).
    
    Args:
        shape: python-pptx shape object containing a table
        content_item: ContentItem with content_type="TABLE" containing:
                      - content: List of lists representing table data (rows x columns)
                      - table_style: Optional TableStyleInfo with header/content cell styles
    
    Returns:
        bool: True if table was successfully updated, False if validation fails or
              an exception occurs during update.
    """
    if not shape.has_table or content_item.content_type != "TABLE" or not content_item.content:
        return False
    
    try:
        table = shape.table
        new_content = content_item.content
        
        # Parse table content - can be list of lists or pipe-separated strings
        table_data = []
        if isinstance(new_content, list):
            for row in new_content:
                if isinstance(row, str):
                    # Parse pipe-separated format: "col1|col2|col3"
                    table_data.append([cell.strip() for cell in row.split('|')])
                elif isinstance(row, list):
                    table_data.append([str(cell) for cell in row])
        else:
            return False
        
        if not table_data:
            return False
        
        new_rows = len(table_data)
        new_cols = max(len(row) for row in table_data)
        current_rows = len(table.rows)
        current_cols = len(table.columns)
        
        # Store the original table width from the shape
        original_table_width = shape.width
        
        # Adjust columns FIRST (before rows, to ensure proper cell structure)
        if new_cols > current_cols:
            for _ in range(new_cols - current_cols):
                _add_table_column(table, current_cols)
        elif new_cols < current_cols:
            for _ in range(current_cols - new_cols):
                _remove_table_column(table)
        
        # Fix column widths to maintain original table width
        if new_cols != current_cols:
            _redistribute_column_widths(table, original_table_width)
        
        # Adjust rows
        if new_rows > current_rows:
            for _ in range(new_rows - current_rows):
                _add_table_row(table, current_rows)
        elif new_rows < current_rows:
            for _ in range(current_rows - new_rows):
                _remove_table_row(table)
        
        # CRITICAL: Apply styling FIRST, before setting text
        # This creates the formatting template that will be preserved
        if content_item.table_style:
            table_style = content_item.table_style
            
            for row_idx in range(len(table.rows)):
                for col_idx in range(len(table.columns)):
                    cell = table.rows[row_idx].cells[col_idx]
                    
                    # First row gets header style, rest get content style
                    if row_idx == 0:
                        apply_cell_style(cell, table_style.header_style, debug=False)
                    else:
                        apply_cell_style(cell, table_style.content_style, debug=False)
        
        # Now fill cells with data WITHOUT destroying formatting
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_value in enumerate(row_data):
                if row_idx < len(table.rows) and col_idx < len(table.columns):
                    cell = table.rows[row_idx].cells[col_idx]
                    text_frame = cell.text_frame
                    
                    # Get the first paragraph (there should be at least one after styling)
                    if not text_frame.paragraphs:
                        continue
                    
                    paragraph = text_frame.paragraphs[0]
                    
                    # Store formatting from existing runs (if any)
                    stored_format = None
                    if paragraph.runs:
                        first_run = paragraph.runs[0]
                        stored_format = {
                            'size': first_run.font.size,
                            'bold': first_run.font.bold,
                            'italic': first_run.font.italic,
                            'name': first_run.font.name,
                        }
                        try:
                            if first_run.font.color and first_run.font.color.rgb:
                                stored_format['color_rgb'] = first_run.font.color.rgb
                        except:
                            pass
                    
                    # If no stored format (new cell), use table_style
                    if not stored_format and content_item.table_style:
                        table_style = content_item.table_style
                        style = table_style.header_style if row_idx == 0 else table_style.content_style
                        stored_format = {
                            'size': Pt(style.font_size),
                            'bold': style.font_bold,
                            'italic': style.font_italic,
                            'name': style.font_name,
                        }
                        if style.font_color_rgb:
                            stored_format['color_rgb'] = style.font_color_rgb
                    
                    # Clear the paragraph (not text_frame, to preserve paragraph-level formatting)
                    paragraph.clear()
                    
                    # Add new run with text
                    run = paragraph.add_run()
                    run.text = str(cell_value)
                    
                    # Reapply stored formatting if we had it
                    if stored_format:
                        if stored_format.get('size'):
                            run.font.size = stored_format['size']
                        if stored_format.get('bold') is not None:
                            run.font.bold = stored_format['bold']
                        if stored_format.get('italic') is not None:
                            run.font.italic = stored_format['italic']
                        if stored_format.get('name'):
                            run.font.name = stored_format['name']
                        if 'color_rgb' in stored_format:
                            from pptx.dml.color import RGBColor
                            run.font.color.rgb = stored_format['color_rgb']
                    
        return True
    
    except Exception as e:
        print(f"Warning: Failed to update table: {e}")
        import traceback
        traceback.print_exc()
        return False


def apply_cell_style(cell, style: CellStyle, debug=False):
    """
    Apply comprehensive styling to a PowerPoint table cell.
    
    Applies font properties (size, bold, italic, name, color) to all text runs
    in the cell and optionally applies a solid fill color to the cell background.
    
    Note: This function only styles existing runs. For newly created cells with no
    runs, styling must be applied when creating the run (see update_table_component).
    
    Args:
        cell: python-pptx table cell object
        style: CellStyle object containing formatting properties to apply
        debug: Unused parameter, kept for backward compatibility
    """
    from pptx.dml.color import RGBColor
    
    try:
        if not cell.text_frame:
            return
        
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(style.font_size)
                run.font.bold = style.font_bold
                run.font.italic = style.font_italic
                run.font.name = style.font_name
                
                if style.font_color_rgb:
                    run.font.color.rgb = RGBColor(*style.font_color_rgb)
        
        if style.has_fill and style.fill_color_rgb:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*style.fill_color_rgb)
    
    except Exception:
        pass


def _add_table_row(table, template_row_idx=0):
    """
    Add a new row to a PowerPoint table by deep-copying a template row.
    
    Creates a new row with identical formatting, borders, and cell properties
    as the template row. The new row is inserted at the end of the table and
    its cell text is cleared.
    
    Args:
        table: python-pptx table object
        template_row_idx: Zero-based index of the row to use as formatting template.
                          Defaults to 0 (first row). If index exceeds table size,
                          uses the last row.
    """
    from copy import deepcopy
    
    # Use the last row or specified template row
    if template_row_idx >= len(table.rows):
        template_row_idx = len(table.rows) - 1
    
    template_row = table._tbl.tr_lst[template_row_idx]
    new_row = deepcopy(template_row)
    
    # Insert after the last row (maintains XML schema order, e.g. before extLst)
    table._tbl.tr_lst[-1].addnext(new_row)
    
    # Clear text from the newly added row using python-pptx API
    new_row_idx = len(table.rows) - 1
    for col_idx in range(len(table.columns)):
        try:
            cell = table.rows[new_row_idx].cells[col_idx]
            if cell.text_frame:
                cell.text = ""
        except:
            pass


def _remove_table_row(table):
    """
    Remove the last row from a PowerPoint table.
    
    Args:
        table: python-pptx table object with at least 2 rows
    """
    if len(table.rows) > 1:
        table._tbl.remove(table._tbl.tr_lst[-1])


def _add_table_column(table, template_col_idx=0):
    """
    Add a new column to a PowerPoint table by deep-copying a template column.
    
    Creates a new column with identical width, formatting, and cell properties
    as the template column. The new column is appended to the right side of the
    table and its cell text is cleared.
    
    Args:
        table: python-pptx table object
        template_col_idx: Zero-based index of the column to use as formatting template.
                          Defaults to 0 (first column). If index exceeds table size,
                          uses the last column.
    """
    from copy import deepcopy
    
    # Use the last column or specified template column
    if template_col_idx >= len(table.columns):
        template_col_idx = len(table.columns) - 1
    
    # Add to grid (copy width from template column)
    template_gridCol = table._tbl.tblGrid.gridCol_lst[template_col_idx]
    new_gridCol = deepcopy(template_gridCol)
    
    # Insert after last gridCol
    table._tbl.tblGrid.gridCol_lst[-1].addnext(new_gridCol)
    
    # Add cell to each row (copy from template column)
    for row_idx, row in enumerate(table._tbl.tr_lst):
        if template_col_idx < len(row.tc_lst):
            template_cell = row.tc_lst[template_col_idx]
            new_cell = deepcopy(template_cell)
            
            # Insert after last cell in the row
            row.tc_lst[-1].addnext(new_cell)
    
    # Clear text from the newly added column using python-pptx API
    new_col_idx = len(table.columns) - 1
    for row_idx in range(len(table.rows)):
        try:
            cell = table.rows[row_idx].cells[new_col_idx]
            if cell.text_frame:
                cell.text = ""
        except:
            pass


def _remove_table_column(table):
    """
    Remove the last column from a PowerPoint table.
    
    Removes the column from both the table grid and all rows.
    
    Args:
        table: python-pptx table object with at least 2 columns
    """
    if len(table.columns) > 1:
        # Remove from grid
        table._tbl.tblGrid.remove(table._tbl.tblGrid.gridCol_lst[-1])
        
        # Remove cell from each row
        for row in table._tbl.tr_lst:
            row.remove(row.tc_lst[-1])


def _redistribute_column_widths(table, target_width):
    """
    Redistribute column widths evenly to maintain a fixed total table width.
    
    Divides the target width equally among all columns, with any remainder
    allocated to the last column to account for integer division rounding.
    This prevents tables from extending beyond slide boundaries when columns
    are added or removed.
    
    Args:
        table: python-pptx table object
        target_width: Desired total table width in EMUs (English Metric Units).
                      Typically obtained from shape.width before modifications.
    """
    num_cols = len(table.columns)
    if num_cols == 0:
        return
    
    # Calculate width per column (evenly distributed)
    width_per_col = target_width // num_cols
    
    # Set width for each column in the grid
    for col_idx, gridCol in enumerate(table._tbl.tblGrid.gridCol_lst):
        # For the last column, use remaining width to account for rounding
        if col_idx == num_cols - 1:
            gridCol.w = target_width - (width_per_col * (num_cols - 1))
        else:
            gridCol.w = width_per_col
    
    # Also set cell widths for all cells in the table
    for row in table.rows:
        for col_idx, cell in enumerate(row.cells):
            if col_idx == num_cols - 1:
                cell.width = target_width - (width_per_col * (num_cols - 1))
            else:
                cell.width = width_per_col



def update_chart_component(shape, content_item) -> bool:
    """
    Update a chart component with new data series.

    Args:
        shape: PPTX shape object with chart
        content_item: ContentItem containing the new chart data

    Returns:
        bool: True if update was successful
    """
    if not shape.has_chart or content_item.content_type != "CHART" or not content_item.content:
        return False

    try:
        chart = shape.chart
        new_content = content_item.content

        series_data = None
        categories = None

        if isinstance(new_content, dict):
            series_data = new_content.get('series', [])
            categories = new_content.get('categories', [])
        elif hasattr(new_content, 'series'):
            series_data = new_content.series
            categories = getattr(new_content, 'categories', [])

        if not series_data:
            return False

        # Create new chart data
        from pptx.chart.data import CategoryChartData
        chart_data = CategoryChartData()

        # Use categories from the stored metadata, or generate defaults
        if categories:
            chart_data.categories = [str(cat) for cat in categories]
        else:
            # Generate default categories based on series length
            max_length = max(len(getattr(series, 'values', []) if hasattr(series, 'values') else series.get('values', [])) for series in series_data)
            chart_data.categories = [f"Category {i+1}" for i in range(max_length)]

        # Add series data
        for series in series_data:
            if isinstance(series, dict):
                name = series.get('name', '')
                values = series.get('values', [])
            else:
                name = getattr(series, 'name', '')
                values = getattr(series, 'values', [])

            if name and values:
                chart_data.add_series(name, values)

        # Replace the chart data
        chart.replace_data(chart_data)
        return True

    except Exception as e:
        print(f"Warning: Failed to update chart: {e}")
        return False


def _copy_font_formatting(target_font, source_font):
    """
    Copy font formatting from source to target font.

    Args:
        target_font: Font object to update
        source_font: Font object to copy from
    """
    target_font.size = source_font.size
    target_font.bold = source_font.bold
    target_font.italic = source_font.italic
    target_font.underline = source_font.underline
    try:
        target_font.color.rgb = source_font.color.rgb
    except AttributeError:
        from pptx.dml.color import RGBColor
        target_font.color.rgb = RGBColor(0, 0, 0)
    target_font.name = source_font.name


def apply_content_to_presentation(
    presentation: Presentation,
    content_mapping: PresentationMapping,
    output_path: str = None
) -> str:
    """
    Apply generated content to presentation slides using UUID mapping.

    Args:
        presentation: Presentation object to update
        content_mapping: PresentationMapping with generated content
        output_path: Path to save the updated presentation (optional)

    Returns:
        Path to the saved presentation
    """
    uuid_to_shape = {}

    for slide_idx, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            # Get UUID from alt_text for all shape types
            element_uuid = get_shape_alt_text(shape)
            
            if element_uuid:
                for slide_content in content_mapping.slides:
                    if slide_content.slide == slide_idx and element_uuid in slide_content.content:
                        content_item = slide_content.content[element_uuid]
                        uuid_to_shape[element_uuid] = (slide_idx, shape, content_item)
                        break

    applied_count = 0
    chart_applied_count = 0
    table_applied_count = 0

    for uuid, (slide_idx, shape, content_item) in uuid_to_shape.items():
        if not content_item.content:
            continue

        if content_item.content_type == "CHART" and shape.has_chart:
            if update_chart_component(shape, content_item):
                chart_applied_count += 1
        
        elif content_item.content_type == "TABLE" and shape.has_table:
            if update_table_component(shape, content_item):
                table_applied_count += 1
        
        elif uuid.endswith("-cell") and shape.has_table:
            if update_single_cell_table(shape, content_item):
                applied_count += 1
        
        elif shape.has_text_frame:
            if update_text_component(shape, content_item):
                applied_count += 1

    print(f"Applied content to {applied_count} text elements, {table_applied_count} tables, and {chart_applied_count} charts")

    if output_path is None:
        output_path = "updated_presentation.pptx"

    presentation.save(output_path)
    print(f"[DEBUG] Updated presentation saved to {output_path}")

    return output_path


def save_structure_to_file(structure_data: List[Dict[str, Any]], output_path: str) -> None:
    """
    Save structure data to JSON file.
    
    Args:
        structure_data: List of slide dictionaries
        output_path: Path to save the JSON file
    """
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(structure_data, f, indent=2, ensure_ascii=False)
    
    print(f"Structure saved to {output_path}")


def save_presentation(presentation: Presentation, output_path: str) -> None:
    """
    Save presentation to file.
    
    Args:
        presentation: Presentation object
        output_path: Path to save the presentation
    """
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")

def clear_all_alt_text(presentation: Presentation) -> None:
    """
    Clear alt text from all shapes in the presentation.
    This removes UUIDs after processing is complete.
    
    Args:
        presentation: Presentation object
    """
    cleared_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if get_shape_alt_text(shape):
                if set_shape_alt_text(shape, ""):
                    cleared_count += 1
    
    print(f"Cleared alt text from {cleared_count} elements")


def process_tag(text: str, tag: str) -> Tuple[bool, str]:
    if tag or f"<{tag.upper()}>" in text:
        return True, text.replace(tag or f"<{tag.upper()}>", "")
    return False, text